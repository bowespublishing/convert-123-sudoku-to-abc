from pptx import Presentation
import os
from os import listdir
import PySimpleGUI as gui
from tkinter import *
from tkinter import messagebox
from tkinter import filedialog

x = 0

def CheckforSudoku(path):
    prscheck = Presentation(path)
    slides = [slide for slide in prscheck.slides]
    for slide in slides:
        for shape in slide.shapes:
            if shape.has_table:
                if 'Sudoku' in shape.name:
                    global x
                    x = x + 1


bp_base64 = b'iVBORw0KGgoAAAANSUhEUgAAANUAAAB9CAAAAAAQZHncAAAAGXRFWHRTb2Z0d2FyZQBBZG9iZSBJbWFnZVJlYWR5ccllPAAAAyNpVFh0WE1MOmNvbS5hZG9iZS54bXAAAAAAADw/eHBhY2tldCBiZWdpbj0i77u/IiBpZD0iVzVNME1wQ2VoaUh6cmVTek5UY3prYzlkIj8+IDx4OnhtcG1ldGEgeG1sbnM6eD0iYWRvYmU6bnM6bWV0YS8iIHg6eG1wdGs9IkFkb2JlIFhNUCBDb3JlIDYuMC1jMDAyIDc5LjE2NDQ2MCwgMjAyMC8wNS8xMi0xNjowNDoxNyAgICAgICAgIj4gPHJkZjpSREYgeG1sbnM6cmRmPSJodHRwOi8vd3d3LnczLm9yZy8xOTk5LzAyLzIyLXJkZi1zeW50YXgtbnMjIj4gPHJkZjpEZXNjcmlwdGlvbiByZGY6YWJvdXQ9IiIgeG1sbnM6eG1wPSJodHRwOi8vbnMuYWRvYmUuY29tL3hhcC8xLjAvIiB4bWxuczp4bXBNTT0iaHR0cDovL25zLmFkb2JlLmNvbS94YXAvMS4wL21tLyIgeG1sbnM6c3RSZWY9Imh0dHA6Ly9ucy5hZG9iZS5jb20veGFwLzEuMC9zVHlwZS9SZXNvdXJjZVJlZiMiIHhtcDpDcmVhdG9yVG9vbD0iQWRvYmUgUGhvdG9zaG9wIDIxLjIgKFdpbmRvd3MpIiB4bXBNTTpJbnN0YW5jZUlEPSJ4bXAuaWlkOjRBRUE2NEJBOTFBNTExRUM4ODNDQjgxRUY0MEFCNDg4IiB4bXBNTTpEb2N1bWVudElEPSJ4bXAuZGlkOjRBRUE2NEJCOTFBNTExRUM4ODNDQjgxRUY0MEFCNDg4Ij4gPHhtcE1NOkRlcml2ZWRGcm9tIHN0UmVmOmluc3RhbmNlSUQ9InhtcC5paWQ6NEFFQTY0Qjg5MUE1MTFFQzg4M0NCODFFRjQwQUI0ODgiIHN0UmVmOmRvY3VtZW50SUQ9InhtcC5kaWQ6NEFFQTY0Qjk5MUE1MTFFQzg4M0NCODFFRjQwQUI0ODgiLz4gPC9yZGY6RGVzY3JpcHRpb24+IDwvcmRmOlJERj4gPC94OnhtcG1ldGE+IDw/eHBhY2tldCBlbmQ9InIiPz4EINMUAAALGUlEQVR42u3bd3gVVRrH8e9NgQiEIErb0GuygBGIdJASCKILAg82EMGChRaWjqssKrAooi7qIwESqqgUEQUEFEEpSlECJKFqgBBApUkSSspv/5iZe29iYkwiycPdmX8yc6blM3PPe87MeQd54oStslW2ylbZKltlq2yVrbJVtspW2SpbZas8VHVl/riHurcLDWnepd+rO9M9RZVQEtfUYpMHqfzrNGkVUtUbcLzuOaqpkpT8w0g/YKXHqN4y51f7QiuPU2kQ+B32ONUS4CuPU60GtrjW/bz1o+jP4q6bSxnpGRnOVZkZGZnOhYyMdKtNuPTth1Frj2a6n+HEhgVLv/m5GFXT4dYka2Fzr3IA3g2nX5IkJYbUC5xhrZwaWG2INb+hRv36ByRJJ4dXAfBtt855/J3dbwEo1/O74lJdCIKB5nzyE65mrO52SVILuNu6OyFQ+VdzYRgEX5OkT25z7vIvc90CX6vkjWJSne0G1ROM+ZQwoO4Ly9e82Rrw/0qSXoZyieYdAPjYmE9vDGMlaaUPtIzevXtRO8Bo+PbcAsEz1y5/uXmZY0Wtej4xPn7/+tGVoN5es/gp4LkrkqS3vKBGkklZZayehLcXTxjze73hG0lHK8BIo0ZFQNlDkjQQQi5IUlrM9aJWOby8vADKvXrFLP0SGGBtMg0YISmtIUQYRaG0bUqtFEnSG9DgqiHoYe6RdheMknStHsws7n5g2ad2m6W9oOJJa5NrTeH2JElDIVSStM+HKRFWI3CPYf6xDN7W7nrPkF6uBtOLS9Vl8oTxT99bFXCMyZCkU7fCM65tZgEfSFoLfkck6T+wdYVxO3TmNtgkabYrlkj7HHgfkNQKah4t1vYqZUUj4FlJWueKBZIU62sILgTCYklqQ6XLCX40SpP0CdROkTQAHjsVb0xHt5SBNZIigQrTThdrK/xbGPCRpHehRLxrm+Tq0FeSHoEnJB0syf1SC9hlxPVnJak9eHs5JwdES9JjAAH9thdn3+JwWeNn9DyUc7vAmXdCB0laDMHXpbdgjvQivGLE9fWS0oPINi2VpIxpZQB4IKn4VOoNZROliRDg9m+kN4ZwSTruj/d+KYyyCdK30EaK8ab6b5LS6kHvqHdc03/NaJMwsRZA/ePFpxoF7JbeBq8Y1zbnKmI2Tl1gjo77c5+ktMb4HdIsa1UrmJDzaa8sawT0LD7VcOAHI1oscm2zHXjVapv6azFESdIYmKM+8Kkk6SHontuJU7qCb2yxqTpA+bNSUnno5dpmDLBDkrTfhwZ6jPKJkvQ1PKxaBF60wn9AQm5n3gZ8UlyqzQ7oJkkPgPdWq/TH2+Eu8yGkKX67guktSboaRPA2Xx411hzzhyG5nfnHEs7OVpGrYmoDqyVpqxfcYT4TXetmBTRJ42BUGaPRkkbgHwHLzVVDgRnZjp5q/n0fvGOKWPWapMyL3430Bx53hY07tktSfEew7oe0GQKoeNZY2AhlqWQ9kPzaELjvy0tSxuX4NydL0tUOI2KuS+mrK0HHoo4W5f8e2rJhZQfAo2b/Nu0RwKftiNHhpYAeKc56Xxd4yGqd6wAPOg91KAigUkjLhlV9+dtFSVvAp05YeF0H3LqnuHq3Pm1d3aT0l26xiktNTnNt/ySwzFp4xu23KenckFLOY1X/SdIq5zFCbxgqR9XFIT27tG4a0rxTr4jog1neNxyZ0OR23xIVW72cJbJtHzg44rzz8X3Q4GG/ZLlG7zzSvkmzzg+9tOE3SdLp6Mc7NQ4JG7L2hj1d5X/0IOOXw0cvZN6M75hu/slW2SpbZatsla2yVbbKVtkqW2WrbJWtunGqDZFz1uRQvDFybl7vZQ9HRkWdK1rVlXOJx+KPnL6S5/7h0CaH4nshJI895xnv74tMda5f81qVy5b08r21ZtjL8YVUzW7UfmeOe0blrFrWtPn6G6LKkvXo++T5wqi+dUCT9D+vOuYP1S7cKFWFzn0e7NuhIkCTpEKoloDrBfWfUH0FlDx8o1RGMs6lZQ1xpUwURHW8EvTJxy/wUjC0TbtRKudod0dgXSHq1f5/v3spP/Xqxykzz+pGqxRX2hoP/YtjYFQRx8CsY43d4S7PU0VA1RQpOTY+LtkqS4mLj73srkqIHj943KKk36syjsbFWgkNmd+/OXTguKjDWVSXVk4YPGnDVeeeP8XFnpAk/RobmyQpccHYwZPWpWb9B6NHPTZmYZJ0/VDckYwCqUZCjVRpHQ4+t8o24TDHdMOhjX55ys9Iehp3ObsquQY8ZZQdDHcYg0e9ki3VXr0XCECIM1001Bohfxt6KnmskZsR7Fa1k540BovKTlK8D7WTC6TqDS1ljOG7qXCpOp4KdrZurc5mV1WDJyVJP1QE/AKc7VcU8P1z1o4lPjUP3dQchtYC6JzczlrvWGyde2c159kG7YAG1wuiSqpgXOzcVaGduH3MF3vWDvEHuqbnokpvDVXnJ5zeP7/ZVOsX6NebMk8sWTk1CAg8nU31AbR8EL/+C1fNuBMoZ/5w4ysDHRbu/v79MAj3pVFmQVTPAWv/UOWghZHksrc2MCsX1Sbw/kaSdD3VWa+oFyNJyfcCk7OploE3gTsk6Xp/YJhRTzu4BtBfh1yCUh6q5OFAmP5QRR0rHOwqBUFXclbNhPrp2WNgwD7zjOWhaTbVCsD3a7NnWgNqp0rSUmC8WyDjzvyoBu49EB+zZmwNoPqRPFTRzj1HABtzVs0w0ySzqEZbC/3B/0RW1Sqgn9v/73NAkrpCTWc38UQ5aJK/3q3DYX7Uc0h/rKroGgbeCbyQs+oLoPeFbCrn90LvOFNtsqg+sNZ/aNaDE/5ul0IKy7cKAK9mc68pD1Vz154pNc20hBxiYFcg0G2YPEsrvNxJzKL6yv2EKyRtwMqRsuJzvlT/iHigb7/x0QecVSF3lXvfoj20z0WV2ATglj7f5qRalYdqi6mah5HmWkDVW9kL/5yqKzTLRaVLo0oDOIZdLbhqOvgdK3LV3dApN5V0fEowwMMFV80C3/i/UPU58Fleqqt1zfymnFVS2tJAM4+3YKqP3a+tdH9hVduAedbCGndVI1czdMDXTFbPTSXtC4AHCqyK84WJroPdVVjV4ZIw3Fp40V3l7XqemGIF6NxVehBaF1iV2RTqXLRKd3kXVpVSGxqYeWanAp3J7eFAf2ubk1WgeUYeqkfNqlcgld7A/K5GUno4hVVpAPCKJOl8R6yTqAvANDPIdQY+zKUfeOaE64r8s+Cqy42Aidck6de+OAqt2gw4nt6VeHB2XUKDrNy5TtCrBPd/fSH11KIg1337vWpKyeEbT148teoOKBVTcJW2+wMNxs+bO7wKze4ptEpDASjtA1WPhMESSVI72DQeHBVr+gN0T85N1RwoXaUMwGwVQqXNgVbHp/6JZ806WghV2lDzaJ2Oqgu8K0lqBt/pxRJmLuS/rmV/Fr5sqtLHVbCSHTfk8EY6B9XHuah0ZnQ1gFovnNcA81OBP6O6vDBqdlxOG+8Y0r5x28HrM6TNkZEHJCnzkznzz0sHJ93TJLTHawm/Hz1I+2je7G2SpAsbZkQ8PnrWzrScRg+OR0bPNZ5nVs+JNNjH51hFks7MjZ7jOnrynnXr96ZK6uEWqjxmpCejMUzyONVub7f+jseo+kHlcx6hSt7pfP0yFfPbvJtftc3RMXLf+Suppz/rBgT94hmqUYBXpVo1AwDqxsozVItrOV8/lHrmZ3mISqlfvNCrdVDjToPmnsjX6MFNMGVm5ntM5OafbJWtslW2ylbZKltlq2yVrbJVtspW2ar/N9X/AHKrCNoDTutcAAAAAElFTkSuQmCC'
bpicon_base64 = b'iVBORw0KGgoAAAANSUhEUgAAALQAAAC0CAQAAACXxM65AAAAAXNSR0IArs4c6QAADTlJREFUeNrtnXtwVNUdxz93d/NCIogib0HFFyCgbVF81AmgouPY+KigWKn1MbZq67uiVqctbdXBTu1ofRU66JShttpMdaytiaAI+GpF0GgFBF8xoCJJINkku3v6B2fPbpbdvffc3b27uXt++8ee5N7sufvJ9/7O77x+1yKdWdQykgnMYiqHMogqjGWzbtrYzDoaaaaFDkQ6pHtbNfOpZzLDCRiGWhajlfU0sJSw3alB6mhFmFdOr1bqCKaCTbaR/Jx72N9IM0cbyDkMZQMd6Q+PZQ1Ro8c8vaKsYWw6RY/hCU5K67ONuTGLMRxLE+2kOI1XjAoL8HqFkcmKDvJzzjdqLoCNJkQjIh7e1dHAvoZKQaydelbsAV3NVoYZIgWzbYwjHMBivsFcUBvGfKwAtdQbFgW2emoDjGSyIVFgm8zIAJMYbkgU2IYzKcBMM3RUcAswM8AUw8EDm2KxjQMNh4LbdouwGdb3wLqtdLMBxgrhpo0Z0Aa0MQPagDagjRnQBrQxA9qANqCNeWAhP32ZKk5nIrvZyYe8ab/O0IB2e3Mu4DYq5E/dbOB5VtJMa0mMm3k8emdxDOMJ08F2NtGd188+i78wIOV3nbzPn1lMWwmg9nCBlCVuE90iblHRLO4TJ4phIpSXT79NZLJ1YrqoLPbiMC8rO0507AWhV7wr/iBmiGDeQLeJnaI3pZYdYpHYv3xAr8iouYhYJobk+OkL5H3yExEU1WK6WCi29KnjmXIBPVpERTbbKObkpOubJOhr1W+qxCniIbFN1dAkDioaaA/j6Kk2Qft4HuNuFTXoW49qcFGRx0tcx0n8W7b4ddxPrf87LIcmlbdyFUPZh4OZyzK2EgWglmu43vUl7ZaY+/59NxuZy0P0ABbf4aaiRbSe3TyL1C3cIY4Xlvp9SEwQ94kedewGl59/gfyEG9McqxELRJcQQoiwmOdvH22JRxXoJWmO14sd8miPmJH0b3D+mi3//ta0R2vEUnn8v2Kkn310iBpZivFCmuMNzOFjACq4i0EuauiU7xVpty50cSdbAZjIuX720UEqZaldAk21RuZLX30yl7iooUuNeKS3j/g+MaCSuxjoX9ABFU+EMwz3CFZyI12AxUIO164hrBSdyVbzHAAHcIfnG3YC3lcVk7pNZ0t4XsYfV2tf2m75XpkRYoQGGQTOTd4C6DfQMaXczNbBbyWKU7VR7LRVNKzgSwBGMMuvoIXDYcI1LAHgECZq1tAua6jMcs4W2RBXcrrHzqMIirY762d0AVWcofn5EdkcVmdBKFgsSzNy6IP6QtGwkxcBOFNbc3HQ2WwtWwAYwtxyBx3lNQDGcrAr0FVZ/0ERGXnAD/zuOiwbpQo+IAxY1LkK8GpsPv8N+S8/wdPNwkVQtGXrEr6gF4BvuQI9wOasD2UajQpP9/2VoKKhU0baR7hyHXaK3i5DPLTjGt8pOibPHaW5vSYOOru1qanag/wJOuoYdIW8rEr2KYCi29klS6M8HJv2EHTcdQRT8zrtZfG4IVAQRXfxtSztbxMK9nNF24MeLDsTway9vHS+3a7DsudK4j66xvZK+rmi7W7YAyVgobm6p9OR60gefvIl6IhDRVuMkGdE1RizjuuotHUInSltgc/CO6egA6pHqAs6DnCIo3gbAh4OLHkIutcx6Ph8eYdCp6No2M/mvLjn78kyMu4L0Nl9dJDxsvQ/hyN+qUq1S0VZpa4o5mfQlk0sMYmhsvSuqy64veuoUef7WtHYRMfzlOd8vUCgh6muS48fQSe+VChrr/C7KghbqVlDtyPQNRwgS1/6E3TEEehjGCVLazVjDqeK3lctNmjxcBV+UVxHKEsMfaYqN2rXEFY9y2w2REUlW7378sUBXZFFbSfI0g5e0q4h7joG2fQ7D3DZ2PrIdYzgG7K0kU0uaog5UPTBculuhPf9CdpJY3iR8q+r+Eq7hvjSnOyKPkW+v5LnzUolCDqT66jlh6rzvdhFQxV1oOigWsawHPwJuttG0RaXKe/5V1e3dVzR2brgs2QU3cPf/Ao6bAP6QL4nS7tY5KqGmFT0vhkHi4JcI0sb1PB/2YE+iQmy9ALvuaohKhUdzDjLchTTZWmNx/tpSwZ0LbfKceQISzRH7VIVnWk6y+I02VB2stK/oLP76F+qwG6VXLqbC+jqDN37ebLur1iLt+bhFqXM4Z3F2Vwp/ep2LkuKuPUsMfWVHvTFHCtL/+Rzj0EXpcOSGt6N41fyZo/we7kI0R3obIoex73Kid0L/lV0NG2lFuNYqZayfMDjOdSQUPTeI96V3KImBJax2XPQRe+CT2W5wtzN7XySQw2xLKCnq+HXr3P6Z/YD0NG9XEeQK1jDNPX7+2jIqYbMih7IE6oz9CZvFgF0qBiKHs23GcThnMe0pInatTyYYw0JH913Dmc/HmWMapLvVus6fA/6ZBqoZEBK/+2BnCOBTIq+krNVebH2vE2/jjr2Y5+9usk/ynlLWgJ0RVJdl7FQgd/EQg9nvovuo9PbdBar+e983agBLuJ36qcoi2gFv4Putb2UGTxlO4OtA9qijgeStiM38sci6bmkFL1H1b/J28LDALN5Kgnzf7jYw3UcJdEYJqyLJt5gKJfLvlyIy+nl2jwM+AT5MXcmzbXs4Ga1XLc45mHOipjMmBEWm0STuFvMEDUyl8cU8bhoU4msrnKZW2mS+FjlZ/qsT76mDnG1CJRPlrAu+bWfE4eIwSnHqsS5okUe/1RMzBF0X9slrshTbr1+kbwqMVAa5kO1RT5x7Gnq2YQARvFkHh+x+AW38ifXI4L9sDFMeOlMa+9e53zeAWACC233Czrzi6s4jweLjrlIoDOvJn2bm+VZF2nvm03dihGlnV9zGqtKIglsSSka4F/8lAgwgEXaT0tLgP6cZdzAVO4ooeTGHjYIH8nG6bWsZ1WJZ+V5izWzhR2lkmQuFVUl9oTwojSGlTZnPSTH187hZJeKjtn2Q33tOsKOQEOTnDodzKVa23kSoIMl9zT5ooCusj1vj5+2mMMhrkAHDGgnoOEtXgaghkc0LtEoWnlOZ64DBEvk2TM5WgN0TH2tQDmDjpuTtFHr+UyW6h1fZMwouq852Uq/RW18ONFxzudY0sb+sgYtNBS9i1WydLSav7YHHTONISTmWCocYVgh34c5zq2UrOiAAQ2WI02/Ks+3mK3to8tc0REtL72LV2VppvZlGh/tOJIGWC3fR/d5PoCzr2WiDseKhjdU6UijaHeVOVP0u2pN9RHaNZV1Y1ihCXo3X8jSYUbR7kA7Syncww5ZGm58tDvPHHIIOr5FbbDmApRQebsOXdC9aq68yjZtYOq9Y3y0BuiISm4Z0sxfXuFhRjsfKDqRji2oCa7CuA53oAOaoCvLW9EVSeGXs55kpxodsYyiC6doNJNLJYZJK4yi9UBHNEFHjaLdKbpXORGhpehKE3XogY5KwFFH2jaKdu064km3nYGOJvnoUPmCtpK+vFMMA5Wv7tVyHQEPk8qXHOjkEbWgw4sbqKKPbi3XgeZjGHwGGk1FV6uFu50OfXTMgO4L1ynoeJ7SnZpxtFG0FugBjJYlZ7vEo0muo8aAdu6jj2CELG3UVvSA8gVdpa3oC9Tlve8QdNSA7pvnyAnoQer53RHe0h4bqTGKdgp6jlpzt8FhY5i8sd8o2qGP3pc5qvys4zq609ZWZqArtar9plra2MvfXYCuMqCdjYrcolZFN9NiFI0WPOd2edID0tdqpOyOGNA6ip7GQuXFYzRoDP9HjetwDrqGXySl2n5BLUjXU3RliYEOlRroaho4Tf3UzqVa84aiZEEXKY7O3Em5PwlzD/e4zupVYRSdLXZentQIwmoecJ0Cwig6gwWZxQpmJ+Woa+I82vuDgkrterJVNZ7bOavPNrdnuDanxx1YBnQyjBpqGcMcrujzkJoOnuSaEkpp0s9AJ3upMxEEGcJoxnFkykjbZm7n6ZLLt9GPQCffzMdxXIYOx5NckkOqqYgBbW9hVvMwz+YEK9F47jag09s7LGAVbTl+SiIpZkv5gg5n6Mvt4D0eYXlebvuN6lPXlBhoy7ukcAexlFP6eOpPeJFVNNOcs5ITHfj1HAZs5IwiPJmiREADHM+pHEYtITbzD14uQOM1iQVs41FPHwpZgqDL1wIGgQFtQBszoA1oA9qYAW1AGzOgDWgD2pg3oLsNBA+sO5C3EUpj2awtUGLDtn61zQHWGQoe2LoAjUV7lmL5WIzGAM1Fezpo+VgrzQFaWG9IFNjW0xKgI8ennhuztwY6LKCarXl8fKCxVNvGOMIBIMyFOayPNZbd2rmQcHxn5ScMZVrJrXT1R7zxMI8h4qAFGziRMYZL3u1VrtvjLeK7zDpo5FjGGjJ5tZe5mE/3FBObsttoYhqjjQPJm9NYwzw+jv8Y7OO2nyfEhJLbC9k/m8CHuV497iuNBamjFWFeOb1aqUtN4JDOUVQzn3omM9xMC2i7i1bW08DSvdcop/fIFrWMZBIzmcKhDDLOxMa6aWMzb9PEO7TQkW7h6P8BBU+G2vkWH04AAAAASUVORK5CYII#'


gui.theme('Black')

choose_powerpoint_column = [

    [gui.Push(),gui.Image(bp_base64),gui.Push(),],

    [gui.Push(),gui.Text("Choose your replacements below for each number"),gui.Push()],
    
    [gui.Push(),gui.Text('1'), gui.InputText(key='-no1-', size=(4, 1)), gui.Text('2'), gui.InputText(key='-no2-', size=(4, 1)), gui.Text('3'), gui.InputText(key='-no3-', size=(4, 1)), gui.Text('4'), gui.InputText(key='-no4-', size=(4, 1)), gui.Text('5'), gui.InputText(key='-no5-', size=(4, 1)), gui.Text('6'), gui.InputText(key='-no6-', size=(4, 1)), gui.Text('7'), gui.InputText(key='-no7-', size=(4, 1)), gui.Text('8'), gui.InputText(key='-no8-', size=(4, 1)), gui.Text('9'), gui.InputText(key='-no9-', size=(4, 1)),gui.Push()],


    [gui.Push(), gui.Text("Choose the PowerPoint file you wish to convert into ABC Sudoku Puzzles below."), gui.Push()],

    [gui.Text("Please note you will need to choose a PowerPoint with Sudoku Puzzles already created by the Puzzle Generator inside it!", font=('Arial', 10, 'bold'))],
    
    [
        gui.Push(),
        gui.Text("Choose your PowerPoint File"),
        gui.Push(),
    ],

    [
        gui.Push(),
        gui.In(size=(25, 1), enable_events=True, key="-IMPORTFILE-"),
        gui.FileBrowse(file_types=(("PowerPoint files", "*.pptx"),)),
        gui.Push(),
    ],

    [
        gui.Push(),
        gui.Text("Choose where you want your ABC Sudoku Puzzles to be save to"),
        gui.Push(),
    ],

    [
        gui.Push(),
        gui.In(size=(25, 1), enable_events=True, key="-EXPORTFILE-"),
        gui.FileSaveAs(file_types=(("PowerPoint files", "*.pptx"),)),
        gui.Push(),
    ],

    [
        gui.Push(),
        gui.Button('Cancel'),
        gui.Button('Ok'),
        gui.Push(),

    ],
]


layout = [
    [
        gui.Column(choose_powerpoint_column),
        
    ]
]

window = gui.Window("Convert Sudoku Puzzles Into ABC Sudoku Puzzles", layout, background_color='#000000', icon=(bpicon_base64))


while True:
    event, values = window.read()
    if event == "Exit" or event == 'Cancel' or event == gui.WIN_CLOSED:
        os._exit(0)
        break
    elif event == 'Ok':
        importfile = values["-IMPORTFILE-"]
        exportfile = values["-EXPORTFILE-"]
        replace1 = values["-no1-"]
        replace2 = values["-no2-"]
        replace3 = values["-no3-"]
        replace4 = values["-no4-"]
        replace5 = values["-no5-"]
        replace6 = values["-no6-"]
        replace7 = values["-no7-"]
        replace8 = values["-no8-"]
        replace9 = values["-no9-"]
        importfile2 = '\\'.join(importfile.split('/'))
        CheckforSudoku(importfile2)
        
        IF = values['-IMPORTFILE-']
        EF = values['-EXPORTFILE-']

        confirm = 'true'
        
        if IF == '':
            gui.Popup('You need to select a PowerPoint file to convert from!')
            confirm = 'false'

        if x == 0:
            gui.Popup('You need to select a PowerPoint file with valid Sudoku Puzzles in!')
            confirm = 'false'         

        if EF == '':
            gui.Popup('You need to choose where you want your ABC Sudoku Puzzles to be saved to!')
            confirm = 'false'

        if confirm == 'true':
            break


window.close()


search1 = '1'
search2 = '2'
search3 = '3'
search4 = '4'
search5 = '5'
search6 = '6'
search7 = '7'
search8 = '8'
search9 = '9'


if __name__ == '__main__':

    prs = Presentation(importfile2)
    slides = [slide for slide in prs.slides]
    for slide in slides:
        for shape in slide.shapes:
            if shape.has_table:
                if 'Sudoku' in shape.name:
                    tbl = shape.table
                    row_count = len(tbl.rows)
                    col_count = len(tbl.columns)
                    for i in range(0, row_count):
                        for j in range(0, col_count):
                            cell = tbl.cell(i,j)
                            paragraphs = cell.text_frame.paragraphs
                            for paragraph in paragraphs:
                                for run in paragraph.runs:
                                    
                                    if(run.text.find(search1))!=-1:
                                        if replace1 != '':
                                            run.text = run.text.replace(search1, replace1)           
                                    elif(run.text.find(search2))!=-1:
                                        if replace2 != '':
                                            run.text = run.text.replace(search2, replace2)
                                    elif(run.text.find(search3))!=-1:
                                        if replace3 != '':
                                            run.text = run.text.replace(search3, replace3)
                                    elif(run.text.find(search4))!=-1:
                                        if replace4 != '':
                                            run.text = run.text.replace(search4, replace4)
                                    elif(run.text.find(search5))!=-1:
                                        if replace5 != '':
                                            run.text = run.text.replace(search5, replace5)
                                    elif(run.text.find(search6))!=-1:
                                        if replace6 != '':
                                            run.text = run.text.replace(search6, replace6)
                                    elif(run.text.find(search7))!=-1:
                                        if replace7 != '':
                                            run.text = run.text.replace(search7, replace7)
                                    elif(run.text.find(search8))!=-1:
                                        if replace8 != '':
                                            run.text = run.text.replace(search8, replace8)
                                    elif(run.text.find(search9))!=-1:
                                        if replace9 != '':
                                            run.text = run.text.replace(search9, replace9)


 
    prs.save(exportfile)

completed_column = [

    [gui.Push(),gui.Image(bp_base64),gui.Push(),],

    [gui.Push(),gui.Text("Conversion completed successfully!"),gui.Push()],
    
    
    [
        gui.Push(),
        gui.Button('Ok'),
        gui.Push(),

    ],
]


layout = [
    [
        gui.Column(completed_column),
        
    ]
]

window = gui.Window("Completed!", layout, background_color='#000000', icon=(bpicon_base64))


while True:
    event, values = window.read()
    if event == "Exit" or event == 'Cancel' or event == gui.WIN_CLOSED:
        os._exit(0)
        break
    elif event == 'Ok':
            break


window.close()
os._exit(0)
